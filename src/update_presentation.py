from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_CANDIDATES = [
    Path.home() / "Downloads" / "P2-template (6).pptx",
    Path.home() / "Downloads" / "P2-template (5).pptx",
]
OUTPUT_DIR = PROJECT_ROOT / "presentation"
OUTPUT_PATH = OUTPUT_DIR / "Employee_Burnout_Risk_Prediction.pptx"
COMPARISON_FIGURE = PROJECT_ROOT / "results" / "figures" / "model_comparison.png"


RESULT_ROWS = [
    ("Logistic Regression", "0.851", "0.858", "0.854"),
    ("SVM", "0.848", "0.848", "0.848"),
    ("Random Forest", "0.813", "0.818", "0.815"),
]


def set_run_style(run, size=16, bold=False, color=RGBColor(30, 30, 30)):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def replace_slide8_content(slide):
    title = slide.shapes[0]
    title.text = "empirical study: results"

    content = slide.shapes[1]
    content.text_frame.clear()
    paragraph = content.text_frame.paragraphs[0]
    paragraph.text = "Best model: Logistic Regression (F1 macro = 0.854)"
    paragraph.font.size = Pt(18)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(38, 38, 38)

    for text in [
        "Evaluation: 80/20 Hold-out split with 200 test records",
        "SVM achieved a very similar performance (F1 macro = 0.848)",
        "The Medium class had more confusion because it is the transition zone",
    ]:
        p = content.text_frame.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(55, 55, 55)


def add_results_table(slide):
    rows = len(RESULT_ROWS) + 1
    cols = 4
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.75),
        Inches(2.55),
        Inches(5.15),
        Inches(1.65),
    )
    table = table_shape.table
    headers = ["Algorithm", "Precision", "Recall", "F1-Score"]
    widths = [2.15, 1.0, 1.0, 1.0]

    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(28, 84, 142)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                set_run_style(run, size=11, bold=True, color=RGBColor(255, 255, 255))

    for row_idx, row in enumerate(RESULT_ROWS, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            if row_idx == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(229, 241, 255)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    set_run_style(run, size=10, bold=(row_idx == 1))


def add_comparison_chart(slide):
    if COMPARISON_FIGURE.exists():
        slide.shapes.add_picture(
            str(COMPARISON_FIGURE),
            Inches(6.15),
            Inches(2.35),
            width=Inches(3.15),
        )


def main():
    template_path = next(
        (candidate for candidate in TEMPLATE_CANDIDATES if candidate.exists()),
        None,
    )
    if template_path is None:
        raise FileNotFoundError("Template not found in Downloads.")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(str(template_path))
    slide8 = presentation.slides[7]

    replace_slide8_content(slide8)
    add_results_table(slide8)
    add_comparison_chart(slide8)

    presentation.save(OUTPUT_PATH)
    print(f"Saved updated presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
